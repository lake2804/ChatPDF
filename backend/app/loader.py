"""
Document loader module supporting multiple file types:
PDF, DOCX, PPTX, TXT, Markdown, and Images
"""
import os
import logging
from pathlib import Path
from typing import List, Tuple, Optional
from langchain_core.documents import Document

logger = logging.getLogger(__name__)

# PDF
PDF_AVAILABLE = False
PDF_ERROR = None
try:
    import fitz  # pymupdf
    import pypdf  # Required by PyPDFLoader
    from langchain_community.document_loaders import PyPDFLoader
    PDF_AVAILABLE = True
    logger.info(f"PDF support enabled. pypdf version: {getattr(pypdf, '__version__', 'unknown')}")
except ImportError as e:
    PDF_ERROR = str(e)
    logger.warning(f"PDF support disabled: {e}")

# DOCX
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PPTX
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Images
from PIL import Image
import base64
import io


def load_pdf(path: str) -> List[Document]:
    """Load PDF file and extract text with page metadata."""
    if not PDF_AVAILABLE:
        error_msg = "PDF support requires both PyMuPDF (fitz) and pypdf. "
        error_msg += "Install with: pip install pymupdf pypdf"
        if PDF_ERROR:
            error_msg += f"\nOriginal error: {PDF_ERROR}"
        raise ImportError(error_msg)
    
    # Double check pypdf is available at runtime
    try:
        import pypdf
        import sys
        pypdf_path = pypdf.__file__ if hasattr(pypdf, '__file__') else 'unknown'
        logger.info(f"Using pypdf from: {pypdf_path}, Python: {sys.executable}")
    except ImportError as import_err:
        import sys
        error_msg = f"`pypdf` package not found in current Python environment.\n"
        error_msg += f"Python executable: {sys.executable}\n"
        error_msg += f"Please install with: pip install pypdf\n"
        error_msg += f"Or activate the correct virtual environment."
        raise ImportError(error_msg)
    
    try:
        # Verify PyPDFLoader can be imported and used
        try:
            from langchain_community.document_loaders import PyPDFLoader
        except ImportError as loader_err:
            import sys
            error_msg = f"PyPDFLoader cannot be imported: {loader_err}\n"
            error_msg += f"Python executable: {sys.executable}\n"
            error_msg += f"This may indicate langchain-community is not properly installed.\n"
            error_msg += f"Try: pip install --upgrade langchain-community"
            raise ImportError(error_msg)
        
        logger.info(f"Loading PDF with PyPDFLoader: {path}")
        loader = PyPDFLoader(path)
        documents = []
        for i, doc in enumerate(loader.load()):
            # Add page number to metadata
            doc.metadata["page"] = i + 1
            doc.metadata["file_type"] = "pdf"
            doc.metadata["source_file"] = os.path.basename(path)
            documents.append(doc)
        logger.info(f"Successfully loaded {len(documents)} pages from PDF")
        return documents
    except ImportError as e:
        error_msg = str(e)
        import sys
        if "pypdf" in error_msg.lower() or "pypdf" in error_msg:
            detailed_msg = f"`pypdf` package not found at runtime.\n"
            detailed_msg += f"Python executable: {sys.executable}\n"
            detailed_msg += f"Please install with: pip install pypdf\n"
            detailed_msg += f"Or restart backend with correct virtual environment."
            raise ImportError(detailed_msg)
        raise ImportError(f"Error loading PDF: {error_msg}")
    except Exception as e:
        error_msg = str(e)
        logger.error(f"Error processing PDF file {path}: {error_msg}", exc_info=True)
        raise ValueError(f"Error processing PDF file: {error_msg}")


def load_docx(path: str) -> List[Document]:
    """Load DOCX file and extract text."""
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx is required for DOCX support. Install with: pip install python-docx")
    
    doc = DocxDocument(path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
    
    # Combine all paragraphs into one document
    content = "\n\n".join(full_text)
    return [Document(
        page_content=content,
        metadata={
            "file_type": "docx",
            "source_file": os.path.basename(path),
            "paragraphs": len(full_text)
        }
    )]


def load_pptx(path: str) -> List[Document]:
    """Load PPTX file and extract text from slides."""
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required for PPTX support. Install with: pip install python-pptx")
    
    prs = Presentation(path)
    documents = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        if slide_text:
            content = "\n".join(slide_text)
            documents.append(Document(
                page_content=content,
                metadata={
                    "file_type": "pptx",
                    "slide_number": slide_num,
                    "source_file": os.path.basename(path)
                }
            ))
    
    return documents


def load_txt(path: str) -> List[Document]:
    """Load plain text file."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    
    return [Document(
        page_content=content,
        metadata={
            "file_type": "txt",
            "source_file": os.path.basename(path)
        }
    )]


def load_markdown(path: str) -> List[Document]:
    """Load Markdown file."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    
    return [Document(
        page_content=content,
        metadata={
            "file_type": "markdown",
            "source_file": os.path.basename(path)
        }
    )]


def load_image(path: str) -> Tuple[bytes, dict]:
    """Load image file and return bytes with metadata."""
    with open(path, "rb") as f:
        image_bytes = f.read()
    
    # Get image info
    img = Image.open(io.BytesIO(image_bytes))
    metadata = {
        "file_type": "image",
        "source_file": os.path.basename(path),
        "format": img.format,
        "size": img.size,
        "mode": img.mode
    }
    
    return image_bytes, metadata


def extract_images_from_pdf(path: str) -> List[Tuple[bytes, dict]]:
    """Extract images from PDF with page metadata."""
    if not PDF_AVAILABLE:
        return []
    
    doc = fitz.open(path)
    images = []
    for page_index in range(len(doc)):
        page = doc[page_index]
        for img_index, img in enumerate(doc.get_page_images(page_index)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            metadata = {
                "page": page_index + 1,
                "image_index": img_index,
                "source_file": os.path.basename(path)
            }
            images.append((image_bytes, metadata))
    doc.close()
    return images


def extract_images_from_pptx(path: str) -> List[Tuple[bytes, dict]]:
    """Extract images from PPTX with slide metadata."""
    if not PPTX_AVAILABLE:
        return []
    
    try:
        prs = Presentation(path)
        images = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape_index, shape in enumerate(slide.shapes):
                # Check if shape has an image
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        metadata = {
                            "slide_number": slide_num,
                            "image_index": shape_index,
                            "source_file": os.path.basename(path),
                            "image_format": shape.image.ext  # e.g., 'png', 'jpeg'
                        }
                        images.append((image_bytes, metadata))
                        logger.info(f"Extracted image from slide {slide_num}, shape {shape_index}")
                    except Exception as e:
                        logger.warning(f"Error extracting image from slide {slide_num}, shape {shape_index}: {e}")
                        continue
        
        logger.info(f"Extracted {len(images)} images from PPTX file")
        return images
    except Exception as e:
        logger.error(f"Error extracting images from PPTX: {e}")
        return []


def load_document(file_path: str) -> Tuple[List[Document], List[Tuple[bytes, dict]]]:
    """
    Main loader function that routes to appropriate loader based on file extension.
    Returns (text_documents, image_list) where image_list contains (bytes, metadata) tuples.
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    
    text_docs = []
    images = []
    
    if ext == ".pdf":
        text_docs = load_pdf(file_path)
        images = extract_images_from_pdf(file_path)
    elif ext == ".docx":
        text_docs = load_docx(file_path)
    elif ext == ".pptx":
        text_docs = load_pptx(file_path)
        images = extract_images_from_pptx(file_path)
    elif ext == ".txt":
        text_docs = load_txt(file_path)
    elif ext in [".md", ".markdown"]:
        text_docs = load_markdown(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"]:
        img_bytes, metadata = load_image(file_path)
        images.append((img_bytes, metadata))
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    
    return text_docs, images

